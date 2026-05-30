import random
import string
import csv
import tempfile
import random
import string
from docx import Document
from pptx import Presentation
from email.mime.text import MIMEText
from odf.opendocument import OpenDocumentText, OpenDocumentPresentation
from odf.draw import Page, Frame, TextBox
from odf.text import P
from odf.style import MasterPage, PageLayout
from pdfdocument.document import PDFDocument
from openpyxl import Workbook


def get_random_letters(file_size):
    if not file_size:
        file_size = 128

    return "".join(random.choices(string.ascii_letters, k=file_size))


def generate_csv(file_size):
    with tempfile.NamedTemporaryFile(mode="w", delete=False, suffix=".csv",
                                     newline="") as tmp:
        writer = csv.writer(tmp)
        writer.writerow(["Data"])
        writer.writerow([get_random_letters(file_size)])
        return tmp.name


def generate_doc(file_size):
    doc = Document()
    doc.add_paragraph(get_random_letters(file_size))
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".doc")
    doc.save(tmp.name)
    return tmp.name


def generate_docx(file_size):
    doc = Document()
    doc.add_paragraph(get_random_letters(file_size))
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".docx")
    doc.save(tmp.name)
    return tmp.name


def generate_dot(file_size):
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".dot")
    with open(tmp.name, 'w') as f:
        f.write("digraph G { " + random.choice(string.ascii_uppercase)
                + " -> " + random.choice(string.ascii_uppercase) + "; }")
    return tmp.name


def generate_eml(file_size):
    msg = MIMEText(get_random_letters(file_size))
    msg["Subject"] = "Sample"
    msg["From"] = "example@example.com"
    msg["To"] = "recipient@example.com"
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".eml")
    with open(tmp.name, 'w') as f:
        f.write(msg.as_string())
    return tmp.name


def generate_htm(file_size):
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".htm")
    with open(tmp.name, 'w') as f:
        f.write(f"<html><body><h1>{get_random_letters(file_size)}"
                "</h1></body></html>")
    return tmp.name


def generate_html(file_size):
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".html")
    with open(tmp.name, 'w') as f:
        f.write(f"<html><body><h1>{get_random_letters(file_size)}"
                "</h1></body></html>")
    return tmp.name


def generate_log(file_size):
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".log")
    with open(tmp.name, 'w') as f:
        f.write(f"{get_random_letters(file_size)}\n")
    return tmp.name


def generate_msg(file_size):
    msg = MIMEText(get_random_letters(file_size))
    msg["Subject"] = "Sample"
    msg["From"] = "example@example.com"
    msg["To"] = "recipient@example.com"
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".msg")
    with open(tmp.name, 'w') as f:
        f.write(msg.as_string())
    return tmp.name


def generate_odg(file_size):
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".odg")
    with open(tmp.name, "w") as f:
        f.write(get_random_letters(file_size))
    return tmp.name


def generate_odp(file_size):
    doc = OpenDocumentPresentation()
    page_layout = PageLayout(name="DefaultLayout")
    doc.automaticstyles.addElement(page_layout)
    master_page = MasterPage(name="Default", pagelayoutname="DefaultLayout")
    doc.masterstyles.addElement(master_page)
    slide = Page(name="Slide1", masterpagename="Default")
    doc.presentation.addElement(slide)
    frame = Frame(width="10cm", height="3cm", x="1cm", y="1cm")
    slide.addElement(frame)
    textbox = TextBox()
    frame.addElement(textbox)
    p = P(text=get_random_letters(file_size))
    textbox.addElement(p)
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".odp")
    doc.save(tmp.name)
    return tmp.name


def generate_ods(file_size):
    doc = OpenDocumentPresentation()
    page_layout = PageLayout(name="Layout1")
    doc.automaticstyles.addElement(page_layout)
    master_page = MasterPage(name="Master1", pagelayoutname="Layout1")
    doc.masterstyles.addElement(master_page)
    slide = Page(masterpagename="Master1", name="Slide1")
    doc.presentation.addElement(slide)
    frame = Frame(width="10cm", height="3cm", x="1cm", y="1cm")
    slide.addElement(frame)
    textbox = TextBox()
    frame.addElement(textbox)
    p = P(text=get_random_letters(file_size))
    textbox.addElement(p)
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".ods")
    doc.save(tmp.name)
    return tmp.name


def generate_odt(file_size):
    doc = OpenDocumentText()
    p = P(text=get_random_letters(file_size))
    doc.text.addElement(p)
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".odt")
    doc.save(tmp.name)
    return tmp.name


def generate_pages(file_size):
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pages")
    with open(tmp.name, "w") as f:
        f.write(get_random_letters(file_size))
    return tmp.name


def generate_pdf(file_size):
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pdf")
    pdf = PDFDocument(tmp.name)
    pdf.init_report()
    pdf.h1(get_random_letters(file_size))
    pdf.generate()
    return tmp.name


def generate_ppt(file_size):
    ppt = Presentation()
    slide = ppt.slides.add_slide(ppt.slide_layouts[0])
    slide.shapes.title.text = get_random_letters(file_size)
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".ppt")
    ppt.save(tmp.name)
    return tmp.name


def generate_pptx(file_size):
    ppt = Presentation()
    slide = ppt.slides.add_slide(ppt.slide_layouts[0])
    slide.shapes.title.text = get_random_letters(file_size)
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    ppt.save(tmp.name)
    return tmp.name


def generate_rtf(file_size):
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".rtf")
    with open(tmp.name, "w") as f:
        f.write(
            "{\\rtf1\\ansi\\deff0{\\fonttbl{\\f0\\fnil Arial;}}\\f0\\fs24 "
            + get_random_letters(file_size) + "}")
    return tmp.name


def generate_txt(file_size):
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".txt")
    with open(tmp.name, "w") as f:
        f.write(get_random_letters(file_size))
    return tmp.name


def generate_xls(file_size):
    wb = Workbook()
    ws = wb.active
    ws.append(["Data"])
    ws.append([get_random_letters(file_size)])
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".xls")
    wb.save(tmp.name)
    return tmp.name


def generate_xlsx(file_size):
    wb = Workbook()
    ws = wb.active
    ws.append(["Data"])
    ws.append([get_random_letters(file_size)])
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".xlsx")
    wb.save(tmp.name)
    return tmp.name


def generate_xmind(file_size):
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".xmind")
    with open(tmp.name, "w") as f:
        f.write(get_random_letters(file_size))
    return tmp.name


FILE_GENERATORS = {
    ".csv": generate_csv,
    ".doc": generate_doc,
    ".docx": generate_docx,
    ".dot": generate_dot,
    ".eml": generate_eml,
    ".htm": generate_htm,
    ".html": generate_html,
    ".log": generate_log,
    ".msg": generate_msg,
    ".odg": generate_odg,
    ".odp": generate_odp,
    ".ods": generate_ods,
    ".odt": generate_odt,
    ".pages": generate_pages,
    ".pdf": generate_pdf,
    ".ppt": generate_ppt,
    ".pptx": generate_pptx,
    ".rtf": generate_rtf,
    ".txt": generate_txt,
    ".xls": generate_xls,
    ".xlsx": generate_xlsx,
    ".xmind": generate_xmind,
}
EXTENSIONS = list(FILE_GENERATORS)


def get_file_with_extension(extension, file_size):
    supported_exts = list(FILE_GENERATORS)
    if extension in supported_exts:
        try:
            return FILE_GENERATORS[extension](file_size)
        except Exception as e:
            print(f"Unable to process {extension} due to {e}")
    raise Exception(f"{extension} is not a supported extension")
